import io
import os
import re
import csv
import logging
import tempfile
from pathlib import Path

# ── Lightweight parsing libraries (graceful import) ─────────────────────────

try:
    import pymupdf4llm
except ImportError:
    pymupdf4llm = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation as PptxPresentation
except ImportError:
    PptxPresentation = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import pdfplumber
except ImportError:
    pdfplumber = None

logger = logging.getLogger(__name__)


class LightweightDocumentParser:
    """
    Lightweight multi-format document parser for resource-constrained environments.

    Supported formats:
    - .pdf  → pymupdf4llm (primary) + pdfplumber (table OCR fallback)
    - .docx → python-docx (paragraphs + tables)
    - .pptx → python-pptx (slide text)
    - .xlsx → openpyxl (all sheets → Markdown tables)
    - .xls  → openpyxl (via compatibility read)

    All outputs are structured as Markdown, compatible with smart_markdown_chunker.
    """

    SUPPORTED_EXTENSIONS = {'.pdf', '.docx', '.pptx', '.xlsx', '.xls', '.csv'}

    # Minimum ratio of plumber-detected table chars vs full page text
    # that triggers a pdfplumber table-enrichment pass.
    _TABLE_DENSITY_THRESHOLD = 0.15

    def parse(self, file_bytes: bytes, filename: str) -> str:
        """
        Parse binary file content to Markdown.

        Args:
            file_bytes: Raw bytes of the document.
            filename:   Original filename (used for extension routing).

        Returns:
            Markdown string ready for smart_markdown_chunker.

        Raises:
            ValueError: Unsupported file extension.
            RuntimeError: Missing required library or extraction failure.
        """
        ext = Path(filename).suffix.lower()
        if ext not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported file type: {ext}")

        dispatch = {
            '.pdf':  self._parse_pdf,
            '.docx': self._parse_docx,
            '.pptx': self._parse_pptx,
            '.xlsx': self._parse_excel,
            '.xls':  self._parse_excel,
            '.csv':  self._parse_csv,
        }
        return dispatch[ext](file_bytes)

    # ── PDF ──────────────────────────────────────────────────────────────────

    def _parse_pdf(self, file_bytes: bytes) -> str:
        """
        Two-pass PDF extraction:
        1. pymupdf4llm fast markdown pass (handles most text + simple tables).
        2. pdfplumber table-enrichment pass (replaces inline tables in output
           when table density on a page exceeds _TABLE_DENSITY_THRESHOLD).
        """
        if not pymupdf4llm:
            raise RuntimeError(
                "pymupdf4llm is not installed. Run: pip install pymupdf4llm"
            )

        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name

        try:
            # Pass 1: fast text + layout extraction
            md_text: str = pymupdf4llm.to_markdown(tmp_path)

            # Pass 2: table enrichment via pdfplumber (only if available)
            if pdfplumber:
                md_text = self._enrich_pdf_tables(tmp_path, md_text)
            else:
                logger.debug(
                    "pdfplumber not installed – PDF table OCR pass skipped. "
                    "Install with: pip install pdfplumber"
                )

            return md_text

        except Exception as e:
            logger.error(f"PDF parse failed: {e}")
            raise RuntimeError(f"PDF parse error: {e}") from e
        finally:
            if os.path.exists(tmp_path):
                os.remove(tmp_path)

    def _enrich_pdf_tables(self, pdf_path: str, base_md: str) -> str:
        """
        Use pdfplumber to extract precise table structures from each page.
        For pages where pdfplumber finds significant tables, append a
        '### Tables (Page N)' section with the Markdown-formatted tables.

        This is additive – we never remove text from the pymupdf4llm output,
        only append table blocks that may have been lost in pure text extraction.
        """
        extra_sections: list[str] = []

        try:
            with pdfplumber.open(pdf_path) as pdf:
                for page_num, page in enumerate(pdf.pages, start=1):
                    tables = page.extract_tables()
                    if not tables:
                        continue

                    # Check table density to avoid false positives on
                    # pages where lines are decorative (e.g., horizontal rules)
                    page_text = page.extract_text() or ""
                    table_char_count = sum(
                        len(str(cell))
                        for tbl in tables
                        for row in tbl
                        for cell in row
                        if cell
                    )
                    if page_text and (table_char_count / max(len(page_text), 1)) < self._TABLE_DENSITY_THRESHOLD:
                        continue

                    for tbl_idx, table in enumerate(tables, start=1):
                        md_table = self._table_to_markdown(table)
                        if md_table:
                            extra_sections.append(
                                f"### Bảng {tbl_idx} (Trang {page_num})\n\n{md_table}"
                            )

        except Exception as e:
            logger.warning(
                f"pdfplumber table enrichment failed (continuing with pymupdf4llm output): {e}"
            )

        if extra_sections:
            separator = "\n\n---\n\n## Bảng Biểu Trích Xuất\n\n"
            return base_md + separator + "\n\n".join(extra_sections)

        return base_md

    # ── DOCX ─────────────────────────────────────────────────────────────────

    def _parse_docx(self, file_bytes: bytes) -> str:
        if not DocxDocument:
            raise RuntimeError(
                "python-docx is not installed. Run: pip install python-docx"
            )
        try:
            doc = DocxDocument(io.BytesIO(file_bytes))
            md_lines: list[str] = []

            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                style_name = para.style.name.lower()
                if 'heading 1' in style_name:
                    md_lines.append(f"# {text}")
                elif 'heading 2' in style_name:
                    md_lines.append(f"## {text}")
                elif 'heading 3' in style_name:
                    md_lines.append(f"### {text}")
                else:
                    md_lines.append(text)

            for table in doc.tables:
                md_table = self._docx_table_to_markdown(table)
                if md_table:
                    md_lines.append("")
                    md_lines.extend(md_table)
                    md_lines.append("")

            return "\n\n".join(md_lines)

        except Exception as e:
            logger.error(f"DOCX parse failed: {e}")
            raise RuntimeError(f"DOCX parse error: {e}") from e

    def _docx_table_to_markdown(self, table) -> list[str]:
        """Convert a python-docx Table object to Markdown table lines."""
        num_rows = len(table.rows)
        num_cols = len(table.columns)
        if num_rows == 0 or num_cols == 0:
            return []

        lines: list[str] = []
        for r in range(num_rows):
            row_data: list[str] = []
            for c in range(num_cols):
                try:
                    cell_text = table.cell(r, c).text.strip()
                    clean = re.sub(r'\s+', ' ', cell_text)
                    row_data.append(clean or "")
                except Exception:
                    row_data.append("")
            lines.append("| " + " | ".join(row_data) + " |")
            if r == 0:
                lines.append("| " + " | ".join(["---"] * num_cols) + " |")
        return lines

    # ── PPTX ─────────────────────────────────────────────────────────────────

    def _parse_pptx(self, file_bytes: bytes) -> str:
        if not PptxPresentation:
            raise RuntimeError(
                "python-pptx is not installed. Run: pip install python-pptx"
            )
        try:
            prs = PptxPresentation(io.BytesIO(file_bytes))
            md_lines: list[str] = []
            for slide_num, slide in enumerate(prs.slides, 1):
                md_lines.append(f"## Slide {slide_num}")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text:
                            md_lines.append(text)
                md_lines.append("\n---")
            return "\n\n".join(md_lines)
        except Exception as e:
            logger.error(f"PPTX parse failed: {e}")
            raise RuntimeError(f"PPTX parse error: {e}") from e

    # ── Excel (.xlsx / .xls) ─────────────────────────────────────────────────

    def _parse_excel(self, file_bytes: bytes) -> str:
        """
        Parse Excel workbooks (.xlsx / .xls) to Markdown.

        Strategy per sheet:
        - Sheet name → ## heading.
        - Each sheet's data → Markdown table (first row = header).
        - Empty sheets are skipped.
        - Merged cells: openpyxl unmerges to top-left value, others blank.
        """
        if not openpyxl:
            raise RuntimeError(
                "openpyxl is not installed. Run: pip install openpyxl"
            )
        try:
            wb = openpyxl.load_workbook(
                io.BytesIO(file_bytes),
                read_only=True,
                data_only=True  # Return computed cell values, not formulas
            )
        except Exception as e:
            logger.error(f"Excel open failed: {e}")
            raise RuntimeError(f"Excel parse error: {e}") from e

        md_sections: list[str] = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(values_only=True))

            # Skip fully empty sheets
            non_empty_rows = [
                r for r in rows if any(cell is not None and str(cell).strip() for cell in r)
            ]
            if not non_empty_rows:
                logger.debug(f"Skipping empty Excel sheet: '{sheet_name}'")
                continue

            md_lines: list[str] = [f"## {sheet_name}"]
            md_lines.append("")

            # First non-empty row = header
            header_row = non_empty_rows[0]
            headers = [str(c) if c is not None else "" for c in header_row]
            md_lines.append("| " + " | ".join(headers) + " |")
            md_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")

            for row in non_empty_rows[1:]:
                # Normalize row length to match header count
                cells: list[str] = []
                for i in range(len(headers)):
                    val = row[i] if i < len(row) else None
                    cell_str = str(val) if val is not None else ""
                    # Clean newlines that would break MD table rows
                    cell_str = re.sub(r'[\r\n]+', ' ', cell_str).strip()
                    cells.append(cell_str)
                md_lines.append("| " + " | ".join(cells) + " |")

            md_sections.append("\n".join(md_lines))

        wb.close()

        if not md_sections:
            logger.warning("Excel file produced no extractable content.")
            return ""

        return "\n\n".join(md_sections)

    # ── CSV ──────────────────────────────────────────────────────────────────

    def _parse_csv(self, file_bytes: bytes) -> str:
        """
        Parse CSV to a Markdown table using the built-in csv module.
        Auto-detects delimiter (comma, semicolon, tab).
        First row = header. No extra dependencies required.
        """
        try:
            text = file_bytes.decode("utf-8", errors="replace")
            dialect = csv.Sniffer().sniff(text[:2048], delimiters=",;\t")
            reader = csv.reader(io.StringIO(text), dialect)
            rows = [r for r in reader if any(c.strip() for c in r)]
        except Exception:
            # Sniffer may fail on very simple files — fallback to comma
            reader = csv.reader(io.StringIO(file_bytes.decode("utf-8", errors="replace")))
            rows = [r for r in reader if any(c.strip() for c in r)]

        if not rows:
            logger.warning("CSV file produced no extractable rows.")
            return ""

        max_cols = max(len(r) for r in rows)
        normalized = [r + [""] * (max_cols - len(r)) for r in rows]

        lines: list[str] = []
        lines.append("| " + " | ".join(normalized[0]) + " |")
        lines.append("| " + " | ".join(["---"] * max_cols) + " |")
        for row in normalized[1:]:
            lines.append("| " + " | ".join(row) + " |")

        return "\n".join(lines)

    # ── Shared Helpers ────────────────────────────────────────────────────────

    @staticmethod
    def _table_to_markdown(table: list[list]) -> str:
        """
        Convert a pdfplumber raw table (list of rows, each a list of cell strings)
        to a Markdown table string.

        Rules:
        - None cells → empty string.
        - Newlines inside cells → replaced with space.
        - First row treated as header.
        - Rows with all-empty cells are skipped (often separator artifacts).
        """
        if not table:
            return ""

        cleaned: list[list[str]] = []
        for row in table:
            cells = [
                re.sub(r'[\r\n]+', ' ', str(c).strip()) if c else ""
                for c in row
            ]
            if any(cells):  # skip empty rows
                cleaned.append(cells)

        if not cleaned:
            return ""

        # Normalize all rows to same column count
        max_cols = max(len(r) for r in cleaned)
        normalized = [r + [""] * (max_cols - len(r)) for r in cleaned]

        lines: list[str] = []
        lines.append("| " + " | ".join(normalized[0]) + " |")
        lines.append("| " + " | ".join(["---"] * max_cols) + " |")
        for row in normalized[1:]:
            lines.append("| " + " | ".join(row) + " |")

        return "\n".join(lines)
